#!/usr/bin/env python3
"""Generate Blade Runner 2049 themed presentation slides using Qwen Image Plus."""

import json
import os
import time
import urllib.request
from pathlib import Path

import dashscope
from dashscope import MultiModalConversation
from pptx import Presentation
from pptx.util import Inches

# ── Config ─────────────────────────────────────────────────────────────
API_KEY = os.environ["QWEN_API_KEY"]
MODEL = "qwen-image-max"
OUTDIR = Path("slides-qwen-max")
OUTPUT = "Agentic_SE_Playbook_QwenMax.pptx"
SIZE = "1344*768"  # closest 16:9 supported

dashscope.base_http_api_url = "https://dashscope.aliyuncs.com/api/v1"

# ── Blade Runner 2049 Design System ───────────────────────────────────
STYLE = """You are a world-class cinematic presentation designer creating a single slide image for a corporate presentation.

VISUAL STYLE — BLADE RUNNER 2049 CINEMATIC:
- Dark, moody, cinematic atmosphere inspired by Blade Runner 2049 (Denis Villeneuve)
- Color palette: deep midnight blue and dark teal backgrounds, warm neon orange as primary accent, cool cyan/teal as secondary accent, muted purple highlights, amber glow effects
- Atmospheric backgrounds: fog, volumetric light shafts, rain particles, distant blurred cityscapes with towering dark buildings, holographic projections, neon reflections on wet surfaces
- Semi-transparent dark panels with thin glowing cyan or orange borders for text readability
- Clean futuristic sans-serif typography — white or amber titles with subtle glow, light gray body text
- Holographic UI elements: thin interface lines, HUD-style corner brackets, scanning lines
- The theme represents AI agents transforming human software engineering — beautiful, awe-inspiring, slightly melancholic

CRITICAL TEXT RULES:
- All text must be EXACTLY as specified — every word, letter, and punctuation mark must be pixel-perfect
- Text must be clearly legible against the dark atmospheric backgrounds
- Use semi-transparent dark overlay panels behind dense text blocks
- Title text should be large, bold, white with subtle glow
- Body text in light gray or white, clean and readable
- No extra text, watermarks, or labels beyond what is specified

"""

# ── 13 Slide Prompts ──────────────────────────────────────────────────
SLIDES = [
    # SLIDE 1: Cover
    {
        "file": "slide-01.png",
        "prompt": STYLE + """SLIDE 1 — COVER

SCENE: A vast Blade Runner 2049 cityscape at twilight. Massive dark skyscrapers disappear into orange-amber fog. A lone silhouette of a software engineer stands on a rain-slicked rooftop, looking out at the city. Above the city, enormous holographic projections flicker — showing code fragments, neural network diagrams, and agent workflow visualizations floating in the misty air. Rain falls gently. Neon reflections shimmer on wet surfaces.

TEXT ON THE SLIDE:
- Small "PLAYBOOK" label in top-left, cyan, uppercase, letter-spaced
- Center, on a wide semi-transparent dark panel:
  Large bold: "Agentic Software Engineering Playbook" (white, amber glow)
  Below: "From OpenAI's Internal Transformation" (light gray)
- Thin cyan horizontal divider line
- Italic: "Software development is undergoing a renaissance in front of our eyes." (pale amber)
- Below: "— Greg Brockman, OpenAI Co-founder" (gray)
- Bottom center: "For Engineering Teams, Tech Leads, CTOs | February 2026" (small gray)
- Bottom right: "1 / 13" (cyan)""",
    },

    # SLIDE 2: Inflection Point
    {
        "file": "slide-02.png",
        "prompt": STYLE + """SLIDE 2 — THE DECEMBER 2025 INFLECTION POINT

SCENE: Dramatic split-screen. LEFT: A dim old-world office with a single desk lamp, dusty keyboard, simple terminal — cold blue-gray monochrome, fading. RIGHT: An explosion of neon — a Blade Runner holographic command center with floating AI agent screens, streaming code in cyan and orange. Between them, a vertical stream of digital transformation energy crackling with electricity. Rain visible through windows.

TEXT ON THE SLIDE:
- Title bar at top: "The December 2025 Inflection Point" (white bold, glow)
- LEFT panel (orange border): "Before December 2025" header in orange, then bullets:
  • Codex for unit tests only
  • Limited to simple tasks
  • Experimental tool
- RIGHT panel (cyan border): "After December 2025" header in cyan, then bullets:
  • Codex writes essentially all code
  • Handles operations and debugging
  • Primary development interface
- Bottom section "Real-World Proof" in amber:
  "Sora Android App — 4 engineers · 18 days · vs 3–6 months" with large orange "~10x faster"
  "C Compiler (100K lines) — 16 parallel agents · ~2,000 sessions · $20K API cost" with cyan "Compiles Linux 6.9 kernel"
- Footer: "Catalyst: GPT-5.2-Codex — long-horizon work, large code changes, leading SWE-Bench Pro" (small gray)
- "2 / 13" bottom right""",
    },

    # SLIDE 3: Adoption Metrics
    {
        "file": "slide-03.png",
        "prompt": STYLE + """SLIDE 3 — ADOPTION METRICS: THE SCALE OF CHANGE

SCENE: Three enormous holographic data pillars of light rising from a dark rain-slicked cityscape into a foggy sky. Each pillar displays a giant glowing number. Smaller data particles and graph fragments float like digital rain around them. The scene has Blade Runner 2049 scale and grandeur. Warm orange and cool cyan light interplay. Flying vehicles in the distance.

TEXT ON THE SLIDE:
- Title: "The Scale of Change" (white bold, top-left)
- Three giant holographic metrics displayed on the pillars:
  PILLAR 1 (orange glow): "1M+" huge bold, below "Developers using Codex", small "and growing"
  PILLAR 2 (cyan glow): "2x" huge bold, below "Usage growth since GPT-5.2", small "nearly doubled"
  PILLAR 3 (purple glow): "20x" huge bold, below "Growth since August 2025", small "explosive adoption"
- Bottom panel: "Enterprise Customers" in amber bold
  "Cisco · Ramp · Virgin Atlantic · Vanta · Duolingo · Gap" (white)
- "3 / 13" bottom right""",
    },

    # SLIDE 4: March 31 Mandate
    {
        "file": "slide-04.png",
        "prompt": STYLE + """SLIDE 4 — THE MARCH 31, 2026 MANDATE

SCENE: A stark Blade Runner-style corporate interior — think Niander Wallace's enormous minimalist office. Vast dark room with thin water line reflecting on the floor. Dramatic shafts of warm amber light cut through darkness from narrow windows. A massive holographic countdown display shows "MARCH 31, 2026" in glowing orange numerals. Authoritative, inevitable atmosphere.

TEXT ON THE SLIDE:
- Title: "The March 31, 2026 Mandate" (white bold, amber glow)
- Subtitle: "OpenAI's Organizational Targets" (light gray)
- Card 1 (orange border): Large "1" in orange + "Tool of First Resort" (white bold)
  "For any technical task, humans interact with an agent rather than using an editor or terminal." (light gray)
- Card 2 (cyan border): Large "2" in cyan + "Safe & Productive Default" (white bold)
  "Agent utilization is explicitly evaluated as safe, but also productive enough that most workflows don't need additional permissions." (light gray)
- Right sidebar panel: "Why March 31?" (amber bold)
  • Top-down organizational transformation
  • Not aspirational — operational deadline
  • Drives cultural change, not just tool adoption
- "4 / 13" bottom right""",
    },

    # SLIDE 5: Traditional vs Agentic
    {
        "file": "slide-05.png",
        "prompt": STYLE + """SLIDE 5 — TRADITIONAL VS AGENTIC ENGINEERING

SCENE: Left side shows a fading ghost of a solitary programmer at a terminal in a dim cubicle — old way, cold blue monochrome, dissolving. Right side shows a vibrant holographic command center with multiple AI agent interfaces orbiting around a standing engineer who orchestrates them — warm orange and cyan neon. Between them, a vertical stream of digital transformation particles flowing left to right. Blade Runner rain in the background.

TEXT ON THE SLIDE:
- Title: "Traditional vs Agentic Engineering" (white bold)
- Comparison table on semi-transparent dark panel with cyan border:
  Headers: "Dimension" (white) | "Traditional" (orange) | → | "Agentic" (cyan)
  Row 1: First Action | Open IDE, write code | Prompt agent with intent
  Row 2: Code Production | Human writes line-by-line | Agent generates, human reviews
  Row 3: Debugging | Manual trace, breakpoint | Agent diagnoses and fixes
  Row 4: Testing | Write tests after code | Agent writes tests first
  Row 5: Documentation | Write docs after shipping | Spec-driven development
  Row 6: Knowledge | Code comments, wiki | AGENTS.md + Skills
  Row 7: Tool Usage | Click UI, manual steps | Agent invokes via CLI/MCP
  Row 8: Speed | Hours/days per feature | Minutes with feedback loops
- "5 / 13" bottom right""",
    },

    # SLIDE 6: 6-Pillar Overview
    {
        "file": "slide-06.png",
        "prompt": STYLE + """SLIDE 6 — THE 6-PILLAR PLAYBOOK

SCENE: Six towering monolithic pillars rising from a reflective dark wet surface, inspired by Blade Runner 2049 massive architecture. Each pillar glows from within — some orange, some cyan, some purple. Thin laser connection lines form a network between them. The ground reflects their light. Holographic UI fragments float around each pillar. Deep orange-amber haze sky above.

TEXT ON THE SLIDE:
- Title: "The 6-Pillar Playbook" (white bold, glow)
- Subtitle: "OpenAI's Implementation Framework" (light gray)
- 6 holographic panel cards in 3x2 grid:
  Card 1 (orange): "1 Adoption" — Try the tools — designate "Agents Captain" per team
  Card 2 (cyan): "2 AGENTS.md" — Create machine-readable project handbook
  Card 3 (amber): "3 Tools" — Make internal tools agent-accessible (CLI/MCP)
  Card 4 (teal): "4 Architecture" — Restructure codebases for agent-first development
  Card 5 (purple): "5 Quality" — "Say no to slop" — maintain review bar
  Card 6 (cyan): "6 Infrastructure" — Build observability, trajectory tracking
- "6 / 13" bottom right""",
    },

    # SLIDE 7: Pillar 1-2
    {
        "file": "slide-07.png",
        "prompt": STYLE + """SLIDE 7 — PILLAR 1-2: ADOPTION & KNOWLEDGE

SCENE: Interior of a Blade Runner corporate war room. Circular table with holographic displays showing network graphs and team structures. Wall-mounted screens display AGENTS.md file contents in cyan monospace text. Silhouetted figures stand watching. Dark room with dramatic amber and cyan side lighting. Rain against floor-to-ceiling windows showing cityscape. Floating file tree structures in the air.

TEXT ON THE SLIDE:
- Title: "Pillar 1-2: Adoption & Knowledge" (white bold)
- LEFT column (orange border panel):
  "1 Cultural Adoption" header
  Engineers "too busy" → "Agents Captain" role per team
  Uncertainty → Shared internal channels
  Lack of exposure → Company-wide Codex hackathon
  Italic quote: "The tools do sell themselves."
- RIGHT column (cyan border panel):
  "2 AGENTS.md & Skills" header
  Code block in dark terminal style:
  project-root/
  +-- AGENTS.md (Project context, conventions)
  +-- skills/ (deploy.sh, test-suite.yml)
  +-- src/
  "Industry Convergence" in amber
  AGENTS.md (OpenAI, Google) · CLAUDE.md (Anthropic) · copilot-instructions.md (GitHub)
- "7 / 13" bottom right""",
    },

    # SLIDE 8: Pillar 3-4
    {
        "file": "slide-08.png",
        "prompt": STYLE + """SLIDE 8 — PILLAR 3-4: TOOLS & ARCHITECTURE

SCENE: A breathtaking Blade Runner engineering bay. A massive holographic blueprint of a software architecture diagram floats in center of a cavernous dark space — interconnected modules, API endpoints, MCP connections as glowing cyan wireframes. Smaller holographic tool interfaces orbit like satellites. Glowing orange circuit lines trace paths on the floor. Steam rises from below. Dramatic volumetric light from above.

TEXT ON THE SLIDE:
- Title: "Pillar 3-4: Tools & Architecture" (white bold)
- LEFT column (orange accent):
  "3 Agent-Accessible Tooling" header
  • Audit tools — Catalog team dependencies
  • Enable access — CLI wrappers or MCP servers
  • Assign ownership — One person per tool integration
  Callout: "MCP (Model Context Protocol)" (cyan bold) — Emerging standard for agent-tool connections
- RIGHT column (cyan accent):
  "4 Agent-First Codebase" header
  • Fast tests — Agents iterate by running tests repeatedly
  • Clean interfaces — Agents reason about one component at a time
  • Strong types — Compiler acts as correctness checker
  • Clear boundaries — Agents work on isolated components
- "8 / 13" bottom right""",
    },

    # SLIDE 9: Pillar 5-6
    {
        "file": "slide-09.png",
        "prompt": STYLE + """SLIDE 9 — PILLAR 5-6: QUALITY & INFRASTRUCTURE

SCENE: A Blade Runner 2049 quality control chamber — like the memory verification lab. Center: a large holographic quality shield glowing amber. Streams of code scroll upward in orange/cyan. On one side, rejected code dissolves into red particles. On the other, approved code solidifies into clean cyan structures. Dark room with concentrated amber spotlights. Monitoring screens showing graphs and dashboards on walls.

TEXT ON THE SLIDE:
- Title: "Pillar 5-6: Quality & Infrastructure" (white bold)
- LEFT column (orange accent):
  '5 "Say No to Slop"' header
  1. Human accountability — Someone accountable for all merged code
  2. Same review bar — Maintain standard for agent-generated code
  3. Author comprehension — No "vibe merging" — understand what you submit
  Warning box: 'The "Slop" Problem' — Functionally correct but poorly maintainable — passes tests but introduces technical debt.
- RIGHT column (cyan accent):
  "6 Infrastructure for Agents" header
  • Observability — Monitor agent performance and success rates
  • Trajectory tracking — Capture reasoning path, not just output
  • Tool management — Governance over agent tool access
  • Feedback loops — Internal usage data drives improvements
- "9 / 13" bottom right""",
    },

    # SLIDE 10: Evolving Engineer Role
    {
        "file": "slide-10.png",
        "prompt": STYLE + """SLIDE 10 — THE EVOLVING ENGINEER ROLE

SCENE: Most iconic Blade Runner visual. A lone engineer stands center of a vast rain-drenched rooftop. Their shadow stretches and transforms — left side is a traditional programmer at a desk, right side morphs into an orchestrator commanding glowing AI agent spheres orbiting around them. City skyline in amber fog. Holographic interfaces float like a constellation. Neon signs flicker. Rain catches light. The transformation from writer to orchestrator.

TEXT ON THE SLIDE:
- Title: "The Evolving Engineer Role" (white bold, glow)
- Subtitle: "From Writer to Orchestrator" (amber)
- Transformation table on dark panel:
  "Old Focus" (dim, fading) → "New Focus" (bright cyan)
  Writing implementation code → Designing agent-implementable systems
  Manual coding → Writing specifications for agents
  Debugging line-by-line → Reviewing and curating agent output
  Ad-hoc knowledge sharing → Maintaining AGENTS.md & skills
  Single-threaded work → Orchestrating multiple agents
  Low-level coding → Higher-order design & communication
- Bottom card: "Spec-Driven Development" (amber bold)
  1. Engineer writes specification → 2. Agent implements against spec → 3. Spec = instruction set + acceptance criteria
- "10 / 13" bottom right""",
    },

    # SLIDE 11: Implementation Roadmap
    {
        "file": "slide-11.png",
        "prompt": STYLE + """SLIDE 11 — IMPLEMENTATION ROADMAP

SCENE: A Blade Runner highway stretching into the distance — a glowing path through dark wasteland. Four towering holographic gateway arches at intervals, each glowing with increasing intensity. First gate dim amber, second brighter orange, third cyan, fourth brilliant white. Road has embedded neon light strips. Fog and dust swirl. Sky transitions from stormy dark (left/present) to clearer amber horizon (right/future). A visual journey through cyberpunk landscape.

TEXT ON THE SLIDE:
- Title: "Implementation Roadmap" (white bold)
- Four phase cards horizontally connected by glowing timeline arrow:
  Phase 1 (amber): "Foundation" Weeks 1-2
  • Designate "Agents Captain"
  • Schedule Codex hackathon
  • Create initial AGENTS.md
  • Audit tooling accessibility
  Phase 2 (orange): "Infrastructure" Weeks 3-4
  • Set up skills repository
  • Create MCP servers (top 5)
  • Establish review guidelines
  • Set up agent observability
  Phase 3 (cyan): "Optimization" Weeks 5-8
  • Optimize tests (<30s target)
  • Refine AGENTS.md from failures
  • Build skills library
  • Define quality metrics
  Phase 4 (white): "Scale" Ongoing
  • Regular AGENTS.md reviews
  • Continuous skills development
  • Infrastructure improvements
- "11 / 13" bottom right""",
    },

    # SLIDE 12: Key Takeaways & Risks
    {
        "file": "slide-12.png",
        "prompt": STYLE + """SLIDE 12 — KEY TAKEAWAYS & OPEN RISKS

SCENE: A vast quiet Blade Runner space. A large holographic eye (replicant iris) fills the upper background — vigilance and awareness. Inside the iris, faint circuit patterns and code fragments visible. Below, two paths diverge — one glowing amber (wisdom/takeaways), one glowing red (risks/warnings). Rain or snow falls gently. Contemplative, powerful — a moment of reckoning.

TEXT ON THE SLIDE:
- Title: "Key Takeaways & Open Risks" (white bold)
- LEFT (55%): "Key Takeaways" in amber bold
  1 Organizational transformation, not tool adoption
  2 Architect role becomes more critical
  3 AGENTS.md is the new README
  4 "Say no to slop" is the hardest pillar
  5 Agent infrastructure is a greenfield opportunity
  6 Navigate deliberately — transition is inevitable
- RIGHT (40%): "Open Risks" in red-orange bold
  Code ownership — Rigorous review, not rubber-stamping
  Technical debt — New metrics for AI-generated drift
  Security — "Safe but productive" balance
  Skill atrophy — Maintain deep understanding for review
- "12 / 13" bottom right""",
    },

    # SLIDE 13: Appendix
    {
        "file": "slide-13.png",
        "prompt": STYLE + """SLIDE 13 — APPENDIX: LANGUAGE-SPECIFIC NOTES

SCENE: A Blade Runner data archive room. Three tall holographic display columns — Rust (orange-tinted), C/C++ (cyan-tinted), Python (purple/green-tinted). Each projects scrolling code. Floating syntax fragments and language symbols drift like embers. Dark room with columns as main light source. Floor reflects their glow. Holographic code brackets {} float in background. Archival, reference-like mood.

TEXT ON THE SLIDE:
- Title: "Appendix: Language-Specific Notes" (white bold)
- Three cards as glowing data terminals:
  Card 1 (orange border): "Rust"
  ✓ Type system = built-in correctness verification
  ⚠ unsafe blocks need explicit AGENTS.md policy
  ⚠ Build times may slow agent iteration
  Card 2 (cyan border): "C / C++"
  ⚠ Less safety guardrails → more prescriptive AGENTS.md
  ⚠ Memory management conventions must be explicit
  Card 3 (purple border): "Python"
  ✓ Large training corpus = agents perform well
  ⚠ Mandate type hints and mypy in AGENTS.md
- Footer: "Source: Greg Brockman's public post on OpenAI's internal agentic transformation | Compiled: February 2026" (small gray)
- "13 / 13" bottom right""",
    },
]


# ── Generate slides ────────────────────────────────────────────────────
def generate_slide(idx: int) -> str | None:
    slide = SLIDES[idx]
    out_path = OUTDIR / slide["file"]

    if out_path.exists() and out_path.stat().st_size > 10000:
        print(f"  [{idx+1}/13] cached: {slide['file']}")
        return str(out_path)

    print(f"  [{idx+1}/13] generating {slide['file']} ...")

    for attempt in range(1, 4):
        try:
            response = MultiModalConversation.call(
                api_key=API_KEY,
                model=MODEL,
                messages=[{"role": "user", "content": [{"text": slide["prompt"]}]}],
                result_format="message",
                stream=False,
                watermark=False,
                prompt_extend=True,
                negative_prompt="blurry text, misspelled words, extra text, watermark, logo",
                size=SIZE,
            )

            if response.status_code == 200:
                content = response.output.choices[0].message.content
                for item in content:
                    if "image" in item:
                        img_url = item["image"]
                        urllib.request.urlretrieve(img_url, str(out_path))
                        size_kb = out_path.stat().st_size / 1024
                        print(f"  [{idx+1}/13] OK: {size_kb:.0f}KB")
                        return str(out_path)

                print(f"  [{idx+1}/13] no image in response, attempt {attempt}/3")
            else:
                print(f"  [{idx+1}/13] HTTP {response.status_code}: {response.message}, attempt {attempt}/3")

        except Exception as e:
            print(f"  [{idx+1}/13] error attempt {attempt}/3: {str(e)[:200]}")

        if attempt < 3:
            wait = attempt * 10
            print(f"  waiting {wait}s ...")
            time.sleep(wait)

    print(f"  [{idx+1}/13] FAILED after 3 attempts")
    return None


# ── Build PPTX ────────────────────────────────────────────────────────
def build_pptx(image_paths: list[str | None]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for i, img_path in enumerate(image_paths):
        slide = prs.slides.add_slide(blank_layout)
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5),
            )
        else:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Slide {i+1} — image generation failed"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xFF, 0x6A, 0x00)

    prs.save(OUTPUT)
    size_mb = Path(OUTPUT).stat().st_size / 1024 / 1024
    print(f"\nPresentation saved: {Path(OUTPUT).resolve()} ({size_mb:.1f}MB)")


# ── Main ───────────────────────────────────────────────────────────────
def main():
    OUTDIR.mkdir(exist_ok=True)
    print(f"Generating 13 slides with {MODEL} (Qwen Image Plus) — Blade Runner 2049 edition...\n")

    image_paths = []
    for i in range(len(SLIDES)):
        result = generate_slide(i)
        image_paths.append(result)
        if i < len(SLIDES) - 1:
            time.sleep(2)

    failed = sum(1 for p in image_paths if p is None)
    if failed:
        print(f"\nWARNING: {failed} slide(s) failed to generate.")

    print("\nBuilding PPTX...")
    build_pptx(image_paths)


if __name__ == "__main__":
    main()
